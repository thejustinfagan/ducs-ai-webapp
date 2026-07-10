#!/usr/bin/env python3
"""Create DUCs.ai GTM PowerPoint presentation"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation
prs = Presentation()

# Color palette: Midnight Executive
NAVY = RGBColor(30, 39, 97)
ICE_BLUE = RGBColor(202, 220, 252)
WHITE = RGBColor(255, 255, 255)
CHARCOAL = RGBColor(54, 69, 79)
CORAL = RGBColor(249, 97, 103)

# SLIDE 1: TITLE
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "DUCs.ai"
title_para.font.size = Pt(64)
title_para.font.bold = True
title_para.font.color.rgb = WHITE
title_para.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.text = "Go-to-Market Strategy"
subtitle_para.font.size = Pt(32)
subtitle_para.font.color.rgb = ICE_BLUE
subtitle_para.alignment = PP_ALIGN.CENTER

tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
tagline_frame = tagline_box.text_frame
tagline_para = tagline_frame.paragraphs[0]
tagline_para.text = "Find your next 10 franchisees in 10 minutes, not 10 weeks"
tagline_para.font.size = Pt(20)
tagline_para.font.italic = True
tagline_para.font.color.rgb = WHITE
tagline_para.alignment = PP_ALIGN.CENTER

date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
date_frame = date_box.text_frame
date_para = date_frame.paragraphs[0]
date_para.text = "July 10, 2026 | Justin Fagan"
date_para.font.size = Pt(14)
date_para.font.color.rgb = ICE_BLUE
date_para.alignment = PP_ALIGN.CENTER

# SLIDE 2: THE OPPORTUNITY
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "The Opportunity"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = NAVY

col_width = Inches(2.8)
col_gap = Inches(0.3)
start_x = Inches(0.5)

# Column 1: Market Size
col1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, Inches(1.5), col_width, Inches(3.5))
col1.fill.solid()
col1.fill.fore_color.rgb = ICE_BLUE
col1.line.fill.background()
col1_frame = col1.text_frame
col1_frame.word_wrap = True
p = col1_frame.paragraphs[0]
p.text = "📊 Total Addressable Market"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
p = col1_frame.add_paragraph()
p.text = "\n18,700 target accounts"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = NAVY
p = col1_frame.add_paragraph()
p.text = "• 2,500 franchise brands (50-500 units)"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = col1_frame.add_paragraph()
p.text = "• 15,000 CRE brokers"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = col1_frame.add_paragraph()
p.text = "• 1,200 franchise law firms"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = col1_frame.add_paragraph()
p.text = "\n$4.2M/year TAM"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY

# Column 2: Problem
col2_x = start_x + col_width + col_gap
col2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col2_x, Inches(1.5), col_width, Inches(3.5))
col2.fill.solid()
col2.fill.fore_color.rgb = RGBColor(250, 240, 240)
col2.line.fill.background()
col2_frame = col2.text_frame
col2_frame.word_wrap = True
p = col2_frame.paragraphs[0]
p.text = "❌ Current Solutions Fail"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = CORAL
p = col2_frame.add_paragraph()
p.text = "\n• LinkedIn: Generic B2B, $100/mo"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = col2_frame.add_paragraph()
p.text = "• CoStar: $5K/yr, poor contact data"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = col2_frame.add_paragraph()
p.text = "• Purchased lists: 40% bounce rate"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = col2_frame.add_paragraph()
p.text = "• Manual research: 20 hrs/week"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL

# Column 3: DUCs Solution
col3_x = col2_x + col_width + col_gap
col3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col3_x, Inches(1.5), col_width, Inches(3.5))
col3.fill.solid()
col3.fill.fore_color.rgb = NAVY
col3.line.fill.background()
col3_frame = col3.text_frame
col3_frame.word_wrap = True
p = col3_frame.paragraphs[0]
p.text = "✅ DUCs.ai Advantage"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = WHITE
p = col3_frame.add_paragraph()
p.text = "\n• 27,709 contacts (live FDD data)"
p.font.size = Pt(12)
p.font.color.rgb = WHITE
p = col3_frame.add_paragraph()
p.text = "• 1,438 direct outreach (validated)"
p.font.size = Pt(12)
p.font.color.rgb = WHITE
p = col3_frame.add_paragraph()
p.text = "• 71% email valid, 33% phone valid"
p.font.size = Pt(12)
p.font.color.rgb = WHITE
p = col3_frame.add_paragraph()
p.text = "• $49/mo (1/10th enterprise cost)"
p.font.size = Pt(12)
p.font.color.rgb = WHITE

# SLIDE 3: TARGET PERSONAS
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Target Customers"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.font.color.rgb = NAVY

# Persona 1
persona1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(9), Inches(1.8))
persona1.fill.solid()
persona1.fill.fore_color.rgb = ICE_BLUE
persona1.line.fill.background()
p1_frame = persona1.text_frame
p1_frame.word_wrap = True
p = p1_frame.paragraphs[0]
p.text = "🎯 PRIMARY: Franchise Development Director"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p = p1_frame.add_paragraph()
p.text = "Emerging brands (50-500 units) | Budget: $500-2,000/mo | Goal: Sign 15-20 franchisees/year"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = p1_frame.add_paragraph()
p.text = "Pain: Stale lead lists ($5K+, 40% bounce), manual LinkedIn prospecting (20 hrs/wk), can't verify contacts"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL

# Persona 2
persona2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.4), Inches(9), Inches(1.8))
persona2.fill.solid()
persona2.fill.fore_color.rgb = RGBColor(245, 250, 255)
persona2.line.fill.background()
p2_frame = persona2.text_frame
p2_frame.word_wrap = True
p = p2_frame.paragraphs[0]
p.text = "🎯 SECONDARY: CRE Broker"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p = p2_frame.add_paragraph()
p.text = "Independent/regional | Budget: $200-500/mo | Goal: Close $20M in transactions"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = p2_frame.add_paragraph()
p.text = "Pain: CoStar contact data is garbage, can't find off-market property owners, manual county research"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL

# Persona 3
persona3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.4), Inches(9), Inches(1.6))
persona3.fill.solid()
persona3.fill.fore_color.rgb = RGBColor(250, 250, 250)
persona3.line.fill.background()
p3_frame = persona3.text_frame
p3_frame.word_wrap = True
p = p3_frame.paragraphs[0]
p.text = "🎯 TERTIARY: Franchise Attorney"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p = p3_frame.add_paragraph()
p.text = "Regional law firms | Budget: $100-300/mo | Goal: Reduce client onboarding time by 80%"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL

# SLIDE 4: COMPETITIVE LANDSCAPE
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Competitive Positioning"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = NAVY

table_y = Inches(1.3)
row_height = Inches(0.7)
col_widths = [Inches(1.8), Inches(1.5), Inches(2.2), Inches(3.5)]

competitors = [
    ("Competitor", "Price/mo", "Weakness", "DUCs Edge"),
    ("LinkedIn Sales Nav", "$100", "Generic B2B", "Niche focus, 1/2 price"),
    ("CoStar", "$5,000+", "Poor contact data", "Contact-first, 1% of price"),
    ("ZoomInfo", "$1,250+", "Overkill, stale data", "Affordable, fresh FDD data"),
    ("Loopnet", "$300+", "Discovery only", "Built for outreach campaigns"),
]

for i, (name, price, weakness, edge) in enumerate(competitors):
    y_pos = table_y + (i * row_height)
    is_header = (i == 0)
    
    # Name
    name_box = slide.shapes.add_textbox(Inches(0.5), y_pos, col_widths[0], row_height)
    name_frame = name_box.text_frame
    name_para = name_frame.paragraphs[0]
    name_para.text = name
    name_para.font.size = Pt(12)
    name_para.font.bold = is_header
    name_para.font.color.rgb = NAVY if is_header else CHARCOAL
    
    # Price
    price_box = slide.shapes.add_textbox(Inches(0.5) + col_widths[0], y_pos, col_widths[1], row_height)
    price_frame = price_box.text_frame
    price_para = price_frame.paragraphs[0]
    price_para.text = price
    price_para.font.size = Pt(12)
    price_para.font.color.rgb = CORAL if not is_header else NAVY
    if not is_header:
        price_para.font.bold = True
    
    # Weakness
    weak_box = slide.shapes.add_textbox(Inches(0.5) + col_widths[0] + col_widths[1], y_pos, col_widths[2], row_height)
    weak_frame = weak_box.text_frame
    weak_para = weak_frame.paragraphs[0]
    weak_para.text = weakness
    weak_para.font.size = Pt(11)
    weak_para.font.color.rgb = CHARCOAL
    weak_frame.word_wrap = True
    
    # Edge (with background)
    if not is_header:
        edge_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(0.5) + col_widths[0] + col_widths[1] + col_widths[2], y_pos, col_widths[3], row_height)
        edge_shape.fill.solid()
        edge_shape.fill.fore_color.rgb = NAVY
        edge_shape.line.fill.background()
    
    edge_box = slide.shapes.add_textbox(Inches(0.5) + col_widths[0] + col_widths[1] + col_widths[2], y_pos, col_widths[3], row_height)
    edge_frame = edge_box.text_frame
    edge_para = edge_frame.paragraphs[0]
    edge_para.text = edge
    edge_para.font.size = Pt(12)
    edge_para.font.bold = True
    edge_para.font.color.rgb = WHITE if not is_header else NAVY

# SLIDE 5: GO-TO-MARKET STRATEGY
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "90-Day Go-to-Market Plan"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = NAVY

# Phase 1
phase1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(1.7))
phase1.fill.solid()
phase1.fill.fore_color.rgb = ICE_BLUE
phase1.line.fill.background()
p1_frame = phase1.text_frame
p1_frame.word_wrap = True
p = p1_frame.paragraphs[0]
p.text = "📍 Phase 1: Outbound Dominance (Days 1-30)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
p = p1_frame.add_paragraph()
p.text = "Goal: 12 paying customers | Budget: $302/mo"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = p1_frame.add_paragraph()
p.text = "• Tier 1 Blitz: 426 contacts, 3-email sequence (35% open, 7% reply)"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL
p = p1_frame.add_paragraph()
p.text = "• LinkedIn: 500 connection requests, DM sequence"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL
p = p1_frame.add_paragraph()
p.text = "• CRE Broker: Email + postcard campaign"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL

# Phase 2
phase2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.2), Inches(9), Inches(1.7))
phase2.fill.solid()
phase2.fill.fore_color.rgb = RGBColor(245, 250, 255)
phase2.line.fill.background()
p2_frame = phase2.text_frame
p2_frame.word_wrap = True
p = p2_frame.paragraphs[0]
p.text = "📍 Phase 2: Product-Led Growth (Days 31-60)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
p = p2_frame.add_paragraph()
p.text = "Goal: 25 trial signups/week | Budget: $400/mo"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = p2_frame.add_paragraph()
p.text = "• Free Tier Launch: 100 searches/mo, convert at 20%"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL
p = p2_frame.add_paragraph()
p.text = "• SEO Content: 4 pillar posts (2,000-3,500 words each)"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL
p = p2_frame.add_paragraph()
p.text = "• Referral Program: \"Give 1 month, get 1 month\""
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL

# Phase 3
phase3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1), Inches(9), Inches(1.6))
phase3.fill.solid()
phase3.fill.fore_color.rgb = RGBColor(240, 245, 250)
phase3.line.fill.background()
p3_frame = phase3.text_frame
p3_frame.word_wrap = True
p = p3_frame.paragraphs[0]
p.text = "📍 Phase 3: Partnership Scaling (Days 61-90)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
p = p3_frame.add_paragraph()
p.text = "Goal: 3 partners, 20% revenue from partners | Budget: $800/mo"
p.font.size = Pt(12)
p.font.color.rgb = CHARCOAL
p = p3_frame.add_paragraph()
p.text = "• Franchise Consultancies: White-label, 20% rev share"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL
p = p3_frame.add_paragraph()
p.text = "• CRE Tech Platforms: API integration, licensing"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL
p = p3_frame.add_paragraph()
p.text = "• Paid Ads Test: LinkedIn + Google ($500)"
p.font.size = Pt(11)
p.font.color.rgb = CHARCOAL

# SLIDE 6: FINANCIAL PROJECTIONS
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "90-Day Financial Projections"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = WHITE

# Conservative
cons_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.2), Inches(3.5))
cons_box.fill.solid()
cons_box.fill.fore_color.rgb = WHITE
cons_box.line.fill.background()
cons_frame = cons_box.text_frame
cons_frame.word_wrap = True
p = cons_frame.paragraphs[0]
p.text = "📊 Conservative Scenario"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p = cons_frame.add_paragraph()
p.text = "\nMonth 1: 10 customers, $490 MRR"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL
p = cons_frame.add_paragraph()
p.text = "Month 2: 23 customers, $1,127 MRR"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL
p = cons_frame.add_paragraph()
p.text = "Month 3: 44 customers, $2,156 MRR"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL
p = cons_frame.add_paragraph()
p.text = "\nQ2 Total: $11,200 revenue"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = NAVY
p = cons_frame.add_paragraph()
p.text = "Ending ARR: $25,872"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL

# Aggressive
agg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.3), Inches(4.2), Inches(3.5))
agg_box.fill.solid()
agg_box.fill.fore_color.rgb = ICE_BLUE
agg_box.line.fill.background()
agg_frame = agg_box.text_frame
agg_frame.word_wrap = True
p = agg_frame.paragraphs[0]
p.text = "🚀 Aggressive Scenario (2x Conversion)"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY
p = agg_frame.add_paragraph()
p.text = "\nMonth 1: 20 customers, $980 MRR"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL
p = agg_frame.add_paragraph()
p.text = "Month 2: 47 customers, $2,303 MRR"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL
p = agg_frame.add_paragraph()
p.text = "Month 3: 91 customers, $4,459 MRR"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL
p = agg_frame.add_paragraph()
p.text = "\nQ2 Total: $23,000 revenue"
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = NAVY
p = agg_frame.add_paragraph()
p.text = "Ending ARR: $53,508"
p.font.size = Pt(13)
p.font.color.rgb = CHARCOAL

# Unit Economics
unit_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9), Inches(1.2))
unit_frame = unit_box.text_frame
unit_para = unit_frame.paragraphs[0]
unit_para.text = "💰 Unit Economics: CAC $38 | LTV $588 | LTV:CAC 15.5:1 (Industry benchmark: 3:1)"
unit_para.font.size = Pt(14)
unit_para.font.color.rgb = WHITE
unit_para.alignment = PP_ALIGN.CENTER

# SLIDE 7: DAY 1 CHECKLIST
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = WHITE

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Day 1 Checklist (Start Tomorrow)"
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = NAVY

# Morning
morning_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(9), Inches(1.9))
morning_box.fill.solid()
morning_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
morning_box.line.fill.background()
m_frame = morning_box.text_frame
m_frame.word_wrap = True
p = m_frame.paragraphs[0]
p.text = "☀️ Morning (2 hours)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
tasks = [
    "Export Tier 1 list (426 contacts) from DUCs dashboard",
    "Purchase 3 domains: tryducs.com, getducs.io, ducsdata.com ($45)",
    "Setup Google Workspace (3 users, $36/mo)",
    "Sign up for Instantly.ai (Starter plan, $97/mo)",
    "Connect domains, configure SPF/DKIM"
]
for task in tasks:
    p = m_frame.add_paragraph()
    p.text = f"✓ {task}"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.space_before = Pt(4)

# Afternoon
afternoon_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.4), Inches(9), Inches(1.9))
afternoon_box.fill.solid()
afternoon_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
afternoon_box.line.fill.background()
a_frame = afternoon_box.text_frame
a_frame.word_wrap = True
p = a_frame.paragraphs[0]
p.text = "🌤️ Afternoon (3 hours)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
tasks = [
    "Draft Email #1, #2, #3 (use templates, customize with markets)",
    "Import Tier 1 list to Instantly, map fields",
    "Setup email sequence (Day 1, 3, 5, 7)",
    "Create landing page: ducs.ai/franchise-developers (Carrd.co)",
    "Block calendar: 2 hrs/day (Days 3-7) for reply handling"
]
for task in tasks:
    p = a_frame.add_paragraph()
    p.text = f"✓ {task}"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.space_before = Pt(4)

# Evening
evening_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(9), Inches(1.3))
evening_box.fill.solid()
evening_box.fill.fore_color.rgb = RGBColor(245, 245, 250)
evening_box.line.fill.background()
e_frame = evening_box.text_frame
e_frame.word_wrap = True
p = e_frame.paragraphs[0]
p.text = "🌙 Evening (1 hour)"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = NAVY
tasks = [
    "Draft 5 LinkedIn posts (schedule for Week 2)",
    "Prepare response templates for common objections",
    "Send first 20 emails (test domain reputation)"
]
for task in tasks:
    p = e_frame.add_paragraph()
    p.text = f"✓ {task}"
    p.font.size = Pt(11)
    p.font.color.rgb = CHARCOAL
    p.space_before = Pt(4)

# SLIDE 8: CALL TO ACTION
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

# Big stat
stat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(1.8))
stat_frame = stat_box.text_frame
stat_para = stat_frame.paragraphs[0]
stat_para.text = "426"
stat_para.font.size = Pt(80)
stat_para.font.bold = True
stat_para.font.color.rgb = CORAL
stat_para.alignment = PP_ALIGN.CENTER

stat_label = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(9), Inches(0.8))
stat_label_frame = stat_label.text_frame
stat_label_para = stat_label_frame.paragraphs[0]
stat_label_para.text = "Prime Contacts Ready for Outreach"
stat_label_para.font.size = Pt(22)
stat_label_para.font.color.rgb = WHITE
stat_label_para.alignment = PP_ALIGN.CENTER

# Action items
action_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(2.0))
action_frame = action_box.text_frame
action_para = action_frame.paragraphs[0]
action_para.text = "🎯 Goal: 50 customers, $2,156 MRR by October 10"
action_para.font.size = Pt(18)
action_para.font.color.rgb = ICE_BLUE
action_para.alignment = PP_ALIGN.CENTER

action_para = action_frame.add_paragraph()
action_para.text = "\n📧 Launch Tier 1 Blitz tomorrow (20 emails Day 1)"
action_para.font.size = Pt(16)
action_para.font.color.rgb = WHITE
action_para.alignment = PP_ALIGN.CENTER

action_para = action_frame.add_paragraph()
action_para.text = "📊 Track: 30%+ open rate, 7%+ reply rate, 10+ demos/week"
action_para.font.size = Pt(16)
action_para.font.color.rgb = WHITE
action_para.alignment = PP_ALIGN.CENTER

# Footer
footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.5))
footer_frame = footer.text_frame
footer_para = footer_frame.paragraphs[0]
footer_para.text = "DUCs.ai Go-to-Market Strategy | July 2026"
footer_para.font.size = Pt(11)
footer_para.font.color.rgb = ICE_BLUE
footer_para.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = "/Users/justinfagan/ducs-ai-webapp/DUCS_GTM_Presentation.pptx"
prs.save(output_path)
print(f"✅ Created: {output_path}")
print(f"   Slides: 8")
print(f"   Size: {round(os.path.getsize(output_path) / 1024, 1)} KB")